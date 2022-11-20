import pptx
from aiogram import types, Dispatcher
from Files import work_with_files
from bot_runner import bot


admins_list = {}


def register_admin_handlers(dp1: Dispatcher):
    dp1.register_message_handler(enable_admin_rights, commands=['admin'])
    dp1.register_message_handler(set_admin_is_free, commands=['read'])


async def enable_admin_rights(message: types.Message):
    admins_list[message.from_user.id] = 'free'
    await message.answer('Вы стали администратором!')
    await try_send_files_somebody()

    # from PowerPoint import pptx_saver
    # import json
    #
    # with open(r'D:\Programming\python\Programms\Bots\Hack_GoCodeHackMSK_2022\bot\result.json', 'r', encoding='utf-8') as file:
    #     res = json.load(file)
    # await pptx_saver.save_pptx(res)

    # a = bot.download_file('AgACAgIAAxkBAAICOWN5F48gJDrDUJT-kdcC93yYQbmoAAJxwzEbvRLJS_pP3nmepAMiAQADAgADcwADKwQ')
    # # image = r'D:\other\обои\12.jpg'
    # presentation = pptx.Presentation(
    #     r'D:\Programming\python\Programms\Bots\Hack_GoCodeHackMSK_2022\Files\OldFiles\test 1321.pptx')
    # for layout in presentation.slides:
    #     for a in layout.shapes:
    #         if a.text == 'img1':
    #             layout.shapes.add_picture(a, a.left, a.top, a.width, a.height)
    #             a.element.getparent().remove(a.element)
    # presentation.save('out.pptx')


async def set_admin_is_free(message: types.Message):
    admin_id = admins_list.get(message.from_user.id, None)
    if admin_id is not None:
        admins_list[admin_id] = 'free'
        await message.answer('Отмечено, как просмотренное')
        await try_send_files_somebody()


async def try_send_files_somebody(admin_id=[]):
    free_admins = admin_id
    for admin_id, is_free in admins_list.items():
        if is_free:
            free_admins.append(admin_id)

    for admin in free_admins:
        file1, file2, file3 = work_with_files.get_new_pair_of_files()
        if file1:
            media = types.MediaGroup()
            media.attach_document(open(file1, 'rb'))
            media.attach_document(open(file2, 'rb'))
            media.attach_document(open(file3, 'rb'))
            await bot.send_media_group(admin, media=media)

            work_with_files.move_files_to_old_folder(file1, file2, file3)
            admins_list[admin] = 'busy'
