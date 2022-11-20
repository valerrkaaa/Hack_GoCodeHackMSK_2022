from Files import work_with_files
import pptx
from pptx.shapes.autoshape import Shape


async def save_pptx(result_array, file_name):
    pptx_path_input = work_with_files.get_full_path('PowerPoint\\pptx_template.pptx')
    pptx_path_output = file_name + '.pptx'

    presentation = pptx.Presentation(pptx_path_input)
    for layout in presentation.slides:
        for e in layout.shapes:
            for result_element in result_array:
                if type(e) == Shape:
                    if e.text == result_element['field_name']:
                        if result_element['content_type'] == 'image':
                            if result_element['content'] != '-':
                                layout.shapes.add_picture(result_element['content'], e.left, e.top, e.width, e.height)

                                e.element.getparent().remove(e.element)
                                # TODO os.remove(result_element['content'])
                        elif result_element['content_type'] == 'text':
                            if e.text.lower().strip().strip('\n').find(result_element['field_name'].strip().strip('\n')) != -1:
                                e.text = e.text.replace(result_element['field_name'], result_element['content'])
    presentation.save(pptx_path_output)
