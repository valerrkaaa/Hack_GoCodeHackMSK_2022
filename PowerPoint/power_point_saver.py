from pptx import Presentation


class Pres:

    def __int__(self, summary):
        self.summary_tags = ['name_project', 'about_project', 'effectiveness', 'goals',
                    'risks', 'problems', 'decision', 'strategy', 'first_stage',
                    'first_stage_description', 'second_stage', 'second_stage_discription',
                    'third_stage', 'third_stage_discription', 'final_stage', 'final_stage_discription',
                    'forecast','picture1','picture2','descript_picture1','descript_picture2','conclusion']
        self.summary = summary
        self.presentation = Presentation('Pattern.pptx')


    def present_build(self):
        for slide in self.presentation.slides:
            print(len(self.presentation.slides))
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        #print (run.text)
                        for i in range(len(self.summary_tags)):  # берем параграф, ищем в нем теги
                            if run.text.find(self.summary_tags[i]) != -1:
                                # заменяем все тэги
                                # print(run.text)
                                run.text = run.text.replace(self.summary_tags[i], '1')

        self.presentation.save('test.pptx')
